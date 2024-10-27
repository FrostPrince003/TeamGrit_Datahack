import fitz  # PyMuPDF
import cv2  # OpenCV for contour detection in images
import numpy as np
import io
from PIL import Image
from pptx import Presentation
from docx import Document
import os

# Extract text from PDF
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text += page.get_text()
    return text

# Extract images from PDF
def extract_images_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    image_paths = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        images = page.get_images(full=True)
        
        for img_index, img in enumerate(images):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            img_ext = base_image["ext"]
            image = Image.open(io.BytesIO(image_bytes))
            
            image_path = f"pdf_image_{page_num+1}_{img_index+1}.{img_ext}"
            image.save(image_path)
            image_paths.append(image_path)
    return image_paths

# Extract text from PPT
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

# Extract images from PPT
def extract_images_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    image_paths = []
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                img_bytes = io.BytesIO(image.blob)
                img = Image.open(img_bytes)
                img_ext = image.ext
                image_path = f"ppt_image_{slide_num+1}_{len(image_paths)+1}.{img_ext}"
                img.save(image_path)
                image_paths.append(image_path)
    return image_paths

# Extract text from DOCX
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text

# Extract images from DOCX
def extract_images_from_docx(docx_path):
    doc = Document(docx_path)
    image_paths = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img = rel.target_part.blob
            img_ext = rel.target_ref.split(".")[-1]
            img_name = f"docx_image_{len(image_paths)+1}.{img_ext}"
            with open(img_name, "wb") as f:
                f.write(img)
            image_paths.append(img_name)
    return image_paths

# Contour detection function for all images
def detect_contours_in_image(image_path):
    image = cv2.imread(image_path)
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blurred, 50, 150)

    contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    # Simply return the count of contours found without cropping
    return len(contours)

# Main function to handle different file types
def process_document(file_path):
    # Get file extension
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    # Initialize list for all extracted image paths
    image_paths = []

    # Based on the file extension, call the appropriate functions
    if file_extension == ".pdf":
        print("Processing PDF...")
        text = extract_text_from_pdf(file_path)
        print(text)
        image_paths = extract_images_from_pdf(file_path)

    elif file_extension == ".pptx":
        print("Processing PPT...")
        text = extract_text_from_ppt(file_path)
        print(text)
        image_paths = extract_images_from_ppt(file_path)

    elif file_extension == ".docx":
        print("Processing DOCX...")
        text = extract_text_from_docx(file_path)
        print(text)
        image_paths = extract_images_from_docx(file_path)

    else:
        print("Unsupported file format!")
        return

    # Detect contours in all extracted images
    for image_path in image_paths:
        print(f"Detecting contours in {image_path}...")
        contour_count = detect_contours_in_image(image_path)
        print(f"Found {contour_count} contours in {image_path}")

# Example usage
if __name__ == "__main__":
    file_path = r"E:\djs_hackathon\InstructionsCreatePDFofE-VerifyManual.pdf"  # Replace with your file path
    process_document(file_path)
