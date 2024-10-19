import fitz  # PyMuPDF for PDFs
import cv2  # OpenCV for contour detection in images
import numpy as np
import io
from PIL import Image
from pptx import Presentation
from docx import Document

# Extract text from PDF
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text += page.get_text()
    return text

# Extract images from PDF without highlights
def extract_images_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    image_count = 0
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        images = page.get_images(full=True)
        
        for img_index, img in enumerate(images):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            img_ext = base_image["ext"]
            image = Image.open(io.BytesIO(image_bytes))
            
            image.save(f"pdf_image_{page_num+1}_{img_index+1}.{img_ext}")  # Save image without green highlights
            image_count += 1
    return image_count

# Extract text from PPT
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

# Extract images from PPT
def extract_images_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    image_count = 0
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                img_bytes = io.BytesIO(image.blob)
                img = Image.open(img_bytes)
                img_ext = image.ext
                img.save(f"ppt_image_{slide_num+1}_{image_count+1}.{img_ext}")
                image_count += 1
    return image_count

# Extract text from DOCX
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text

# Extract images from DOCX
def extract_images_from_docx(docx_path):
    doc = Document(docx_path)
    image_count = 0
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img = rel.target_part.blob
            img_ext = rel.target_ref.split(".")[-1]
            img_name = f"docx_image_{image_count+1}.{img_ext}"
            with open(img_name, "wb") as f:
                f.write(img)
            image_count += 1
    return image_count

# Contour detection function for all images
def detect_contours_in_image(image_path):
    image = cv2.imread(image_path)
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blurred, 50, 150)

    contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    
    for contour in contours:
        x, y, w, h = cv2.boundingRect(contour)
        cropped_image = image[y:y + h, x:x + w]  # Extract only the detected image region
        output_image_path = f"cropped_{image_path.split('.')[0]}.jpg"
        cv2.imwrite(output_image_path, cropped_image)
    
    return len(contours)

# Example usage
if __name__ == "__main__":
# Example usage
    pdf_path = r"E:\djs_hackathon\InstructionsCreatePDFofE-VerifyManual.pdf"
    ppt_path = r"C:\Users\ASUS\Downloads\Resume shortlisting PPT (1) (1).pptx"
    docx_path = r"C:\Users\ASUS\Downloads\DL_Report[1][1].docx"
    
    # PDF Extraction
    print("Extracting text from PDF...")
    pdf_text = extract_text_from_pdf(pdf_path)
    print(pdf_text)
    
    print("Extracting images from PDF...")
    pdf_image_count = extract_images_from_pdf(pdf_path)
    print(f"Extracted {pdf_image_count} images from PDF")
    
    # PPT Extraction
    print("Extracting text from PPT...")
    ppt_text = extract_text_from_ppt(ppt_path)
    print(ppt_text)
    
    print("Extracting images from PPT...")
    ppt_image_count = extract_images_from_ppt(ppt_path)
    print(f"Extracted {ppt_image_count} images from PPT")
    
    # DOCX Extraction
    print("Extracting text from DOCX...")
    docx_text = extract_text_from_docx(docx_path)
    print(docx_text)
    
    print("Extracting images from DOCX...")
    docx_image_count = extract_images_from_docx(docx_path)
    print(f"Extracted {docx_image_count} images from DOCX")
    
    # Detect contours in extracted images (example with PDF)
    for i in range(1, pdf_image_count + 1):
        img_path = f"pdf_image_1_{i}.png"  # Update extension if needed
        print(f"Detecting contours in {img_path}...")
        contour_count = detect_contours_in_image(img_path)
        print(f"Found {contour_count} contours in {img_path}")
