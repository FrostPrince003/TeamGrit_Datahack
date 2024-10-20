import os
import io
import fitz  # PyMuPDF for PDF processing
import cv2  # OpenCV for contour detection in images
import numpy as np
from PIL import Image
from pptx import Presentation  # python-pptx for PPTX processing
from docx import Document  # python-docx for DOCX processing
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from youtube_transcript_api import YouTubeTranscriptApi  # YouTube transcript extraction
import google.generativeai as genai  # Google Gemini AI
from dotenv import load_dotenv  # dotenv for environment variable management
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
# Load environment variables
load_dotenv()

# Ensure API key is set correctly from environment variables
api_key = os.getenv("YOUR_API_KEY_HERE")
genai.configure(api_key=api_key)  # Configures the Google Gemini API

# FastAPI app instance
app = FastAPI()

origins = [
    "http://localhost:5173",  # React frontend, for example
    "https://myfrontendapp.com",  # Production frontend
]

# Add CORS middleware to the FastAPI app
app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,  # List of allowed origins
    allow_credentials=True,  # Allows cookies and credentials
    allow_methods=["*"],  # Allows all methods (GET, POST, PUT, etc.)
    allow_headers=["*"],  # Allows all headers
)
model = genai.GenerativeModel("gemini-1.5-flash")
# Default prompt template for quiz generation
prompt_template = """
You are a quiz agent. Generate a set of three questions from the text provided, one easy, one medium, and one hard.
Return the result in JSON format.
"""

# Function to extract transcript from YouTube video using YouTubeTranscriptApi
def extract_transcript(youtube_video_url):
    try:
        video_id = youtube_video_url.split("=")[1]
        transcript_data = YouTubeTranscriptApi.get_transcript(video_id)
        transcript = " ".join([item["text"] for item in transcript_data])
        return transcript
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting transcript: {str(e)}")

# Function to generate content using Google Gemini AI
def generate_gemini_content(transcript_text, prompt):
    try:
        
        response = model.generate_content(prompt + transcript_text)
        return response.text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating content: {str(e)}")

# Route to process YouTube link and return detailed notes
@app.post("/generate_notes/")
async def generate_notes(youtube_video: BaseModel):
    try:
        transcript_text = extract_transcript(youtube_video.youtube_url)
        if not transcript_text:
            raise HTTPException(status_code=404, detail="Transcript not found")
        summary = generate_gemini_content(transcript_text, prompt_template)
        return {"detailed_notes": summary}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Function to extract text from PDF
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text += page.get_text()
    return text

# Function to extract images from PDF
def extract_images_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    image_paths = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        images = page.get_images(full=True)
        for img_index, img in enumerate(images):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            img_ext = base_image["ext"]
            image = Image.open(io.BytesIO(image_bytes))
            image_path = f"pdf_image_{page_num + 1}_{img_index + 1}.{img_ext}"
            image.save(image_path)
            image_paths.append(image_path)
    return image_paths

# Function to extract text from PPT
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

# Function to extract text from DOCX
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text

# Function to detect contours in images (optional)
def detect_contours_in_image(image_path):
    image = cv2.imread(image_path)
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blurred, 50, 150)
    contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    return len(contours)  # Return the count of contours found

# Function to generate questions using Google Gemini AI, with options for quiz or flashcards
def generate_questions_from_text(extracted_text, question_type="quiz"):
    model = genai.GenerativeModel("gemini-pro")
    subject_prompt = "Which subject does this text belong to? Answer in just one word: "
    init_prompt = extracted_text + subject_prompt
    init_response = model.generate_content(init_prompt)
    subject = init_response.text.strip()

    if question_type == "quiz":
        prompt = (
            f"Assume you are a {subject} teacher. Based on the following text, generate three multiple-choice questions "
            "(1 easy, 1 medium, and 1 hard) with four answer choices each. The questions can be of any type like fill "
            "in the blanks, true/false, or multiple-choice. Indicate the correct answer. Format everything in JSON.\n\n"
        )
    else:
        prompt = (
            f"Assume you are a {subject} teacher. Based on the following text, generate three open-ended questions "
            "(1 easy, 1 medium, and 1 hard). Indicate the correct answers. Format everything in JSON.\n\n"
        )

    response = model.generate_content(prompt + extracted_text)
    return response.text

# Main function to process the document and extract text based on file type
def process_document(file_path):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    extracted_text = ""

    if file_extension == ".pdf":
        extracted_text = extract_text_from_pdf(file_path)
    elif file_extension == ".pptx":
        extracted_text = extract_text_from_ppt(file_path)
    elif file_extension == ".docx":
        extracted_text = extract_text_from_docx(file_path)
    elif file_extension in [".png", ".jpg", ".jpeg"]:
        # Directly send image to Gemini; you may need a separate image processing function
        return generate_questions_from_text("Image received, no text extracted.")
    else:
        raise ValueError("Unsupported file format!")

    return extracted_text

# FastAPI endpoint to handle file uploads and return quiz questions
@app.post("/upload/document/")
async def upload_and_generate_questions(file: UploadFile = File(...), question_type: str = "quiz"):
    file_location = f"temp/{file.filename}"
    os.makedirs("temp", exist_ok=True)
    with open(file_location, "wb") as f:
        f.write(await file.read())

    try:
        extracted_text = process_document(file_location)
        questions = generate_questions_from_text(extracted_text, question_type=question_type)
        return JSONResponse(content={"generated_questions": questions})
    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=400)
    finally:
        os.remove(file_location)

@app.post("/upload/image/")
async def image_flashcard(image_path: str = "E:\\djs_hackathon\\TeamGrit_Datahack\\docx_image_2.png"):
    if not os.path.exists(image_path):
        raise HTTPException(status_code=404, detail="Image not found.")

    try:
        # Prepare the prompt based on the image
        prompt = (
            f"Assume you are a teacher. Based on the following image, generate one open-ended question "
            "and indicate the correct answer. Format everything in JSON.\n\n"
        )

        # Generate question using the model
        response = model.generate_content(prompt)
        generated_question = response.text
        
        # Return the generated question and the image URL
        return {
            "generated_question": generated_question,
            "image_url": f"http://192.168.254.146:5000/get_image/?image_path={image_path}"  # Direct link to the image
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/get_image/")
async def get_image(image_path: str):
    # Ensure the path is correctly formatted for your OS
    if not os.path.exists(image_path):
        raise HTTPException(status_code=404, detail="Image not found.")
    
    return FileResponse(image_path)


# Endpoint for quiz generation
@app.post("/upload/quiz/")
async def upload_quiz(file: UploadFile = File(...)):
    return await upload_and_generate_questions(file, question_type="quiz")

# Endpoint for flashcard generation
@app.post("/upload/flashcard/")
async def upload_flashcard(file: UploadFile = File(...)):
    return await upload_and_generate_questions(file, question_type="flashcard")

# Simple GET request to verify server is running
@app.get("/get")
def check_status():
    return "Server is running."

# Run the application
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="192.168.254.146", port=5000)
